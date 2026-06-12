"""
VaultRAG – Document Processor

Extracts plain text from every supported file type.
Returns a string (may be empty if extraction fails).
"""

import io
import json
import logging
import os
from pathlib import Path
from typing import Optional

from config.settings import APP_NAME, TESSERACT_CMD

logger = logging.getLogger(APP_NAME)

# ─── Optional dependency guards ───────────────────────────────────────────────

def _try_import(name):
    try:
        import importlib
        return importlib.import_module(name)
    except ImportError:
        return None


# ─── Per-extension extractors ─────────────────────────────────────────────────

def _extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        parts = []
        for i, page in enumerate(reader.pages):
            try:
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(f"[Page {i+1}]\n{text}")
            except Exception as e:
                logger.warning("PDF page %d extraction failed: %s", i + 1, e)
        return "\n\n".join(parts)
    except Exception as e:
        logger.error("PDF extraction error for %s: %s", path.name, e)
        return ""


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append(" | ".join(c.text.strip() for c in row.cells))
            parts.append("\n".join(rows))
        return "\n\n".join(parts)
    except Exception as e:
        logger.error("DOCX extraction error: %s", e)
        return ""


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                parts.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(parts)
    except Exception as e:
        logger.error("PPTX extraction error: %s", e)
        return ""


def _extract_xlsx(path: Path) -> str:
    try:
        import pandas as pd
        xl = pd.ExcelFile(str(path))
        parts = []
        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            parts.append(f"[Sheet: {sheet}]\n{df.to_markdown(index=False)}")
        return "\n\n".join(parts)
    except Exception as e:
        logger.error("XLSX extraction error: %s", e)
        return ""


def _extract_csv(path: Path) -> str:
    try:
        import pandas as pd
        df = pd.read_csv(str(path))
        return df.to_markdown(index=False)
    except Exception as e:
        logger.error("CSV extraction error: %s", e)
        return ""


def _extract_html(path: Path) -> str:
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(path.read_text(errors="replace"), "html.parser")
        return soup.get_text(separator="\n")
    except Exception as e:
        logger.error("HTML extraction error: %s", e)
        return path.read_text(errors="replace")


def _extract_json(path: Path) -> str:
    try:
        data = json.loads(path.read_text(errors="replace"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception:
        return path.read_text(errors="replace")


def _extract_image(path: Path) -> str:
    if TESSERACT_CMD:
        import pytesseract
        pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(str(path))
        return pytesseract.image_to_string(img)
    except Exception as e:
        logger.warning("OCR extraction failed for %s: %s", path.name, e)
        return ""


# ─── Public API ───────────────────────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {
    ".txt", ".md",
    ".pdf",
    ".docx", ".doc",
    ".pptx",
    ".xlsx", ".xls", ".csv",
    ".html", ".htm",
    ".json",
    ".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif",
}


def extract_text(path: Path) -> Optional[str]:
    """
    Extract all readable text from *path*.
    Returns None if the file type is unsupported.
    Returns "" if supported but extraction yielded nothing.
    """
    ext = path.suffix.lower()

    if ext in (".txt", ".md"):
        return path.read_text(errors="replace")
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext in (".docx", ".doc"):
        return _extract_docx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext in (".xlsx", ".xls"):
        return _extract_xlsx(path)
    if ext == ".csv":
        return _extract_csv(path)
    if ext in (".html", ".htm"):
        return _extract_html(path)
    if ext == ".json":
        return _extract_json(path)
    if ext in (".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"):
        return _extract_image(path)

    logger.warning("Unsupported file type: %s", ext)
    return None
